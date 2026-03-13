#!/usr/bin/env python3
"""
generate-report.py — SunExpress branded PPTX rapor üretici

Kullanım:
    python scripts/generate-report.py --week 2026-03-02
    python scripts/generate-report.py --week 2026-03-02 --output custom.pptx
"""

import argparse
import io
import json
from datetime import datetime
from matplotlib.lines import Line2D
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Branding constants  (templates/branding-spec.md)
# ---------------------------------------------------------------------------
NAVY    = RGBColor(0x1B, 0x3A, 0x5C)
ORANGE  = RGBColor(0xE8, 0x5D, 0x26)
TEAL    = RGBColor(0x2A, 0xBF, 0xBF)
LGRAY   = RGBColor(0xF0, 0xF0, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY   = RGBColor(0x33, 0x33, 0x33)
WARN_BG = RGBColor(0xFF, 0xF3, 0xCD)
CRIT_BG = RGBColor(0xF8, 0xD7, 0xDA)
GREEN   = RGBColor(0x28, 0xA7, 0x45)
RED     = RGBColor(0xDC, 0x35, 0x45)
SGRAY   = RGBColor(0x6C, 0x75, 0x7D)

M_NAVY   = "#1B3A5C"
M_ORANGE = "#E85D26"
M_GRID   = "#E0E0E0"
M_DGRAY  = "#333333"

BASE_DIR = Path(__file__).resolve().parent.parent


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def set_cell_bg(cell, rgb: RGBColor):
    """Set table cell background via raw XML."""
    from lxml import etree
    from pptx.oxml.ns import qn
    tc   = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))
    for sf in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(sf)
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", "%02X%02X%02X" % (rgb[0], rgb[1], rgb[2]))


def set_slide_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def new_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_footer(slide):
    """loodos teal left + SunExpress Airlines navy right."""
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(2.5), Inches(0.32))
    run = txb.text_frame.paragraphs[0].add_run()
    run.text = "loodos"
    run.font.size = Pt(10); run.font.color.rgb = TEAL
    run.font.bold = True;   run.font.name = "Calibri"

    txb2 = slide.shapes.add_textbox(Inches(9.8), Inches(7.1), Inches(3.2), Inches(0.32))
    p2   = txb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = "SunExpress Airlines"
    run2.font.size = Pt(10); run2.font.color.rgb = NAVY; run2.font.name = "Calibri"


def add_title(slide, parts):
    """parts: [(text, bold, RGBColor, pt_size), ...]"""
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.75))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    for text, bold, color, size in parts:
        run = p.add_run()
        run.text = text; run.font.bold = bold
        run.font.color.rgb = color; run.font.size = Pt(size); run.font.name = "Calibri"


# ---------------------------------------------------------------------------
# Table helpers
# ---------------------------------------------------------------------------

def style_header_row(tbl, headers, col_widths=None):
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True; run.font.size = Pt(13)
        run.font.color.rgb = WHITE; run.font.name = "Calibri"
        set_cell_bg(cell, NAVY)
    tbl.rows[0].height = Inches(0.40)   # FIX 4: compact header


def style_data_cell(cell, text, flag=None, align=PP_ALIGN.CENTER, font_size=11):
    cell.text = text
    p   = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size); run.font.name = "Calibri"
    run.font.bold = False          # FIX 8: never bold on anomaly, only bg colour
    if flag == "critical":
        set_cell_bg(cell, CRIT_BG)
    elif flag == "warning":
        set_cell_bg(cell, WARN_BG)
    else:
        set_cell_bg(cell, WHITE)


# ---------------------------------------------------------------------------
# Data / formatting helpers
# ---------------------------------------------------------------------------

def fmt_k(val: float) -> str:
    """FIX 2: K-format for large numbers. 21900→21.9K, 74000→74K, 1370→1.4K"""
    k = val / 1000
    if k >= 10:
        return f"{k:.0f}K"
    else:
        return f"{k:.1f}K"


def fmt_date_range(start: str, end: str) -> str:
    """FIX 1: '2026-01-19','2026-01-25' → '19 Jan - 25 Jan'"""
    ds = datetime.strptime(start, "%Y-%m-%d")
    de = datetime.strptime(end,   "%Y-%m-%d")
    return f"{ds.day} {ds.strftime('%b')} - {de.day} {de.strftime('%b')}"


def fmt_date_short(date_str: str) -> str:
    """FIX 1: '2026-01-19' → '19 Jan'"""
    dt = datetime.strptime(date_str, "%Y-%m-%d")
    return f"{dt.day} {dt.strftime('%b')}"


def fmt_metric(val, col: str) -> str:
    """FIX 2 & 3: K format for big numbers; no % for crash_free."""
    if "rating" in col:
        return f"{val:.1f}"
    elif "crash_free" in col:
        return f"{val:.2f}"          # FIX 2: no % sign
    else:
        return fmt_k(val)            # FIX 2: K format


def build_anomaly_map(anomalies_data: dict) -> dict:
    amap = {}
    for anom in anomalies_data.get("anomalies", []):
        col = anom["column"]
        for wk in anom.get("weeks", []):
            amap[(wk["week_start"], col)] = wk["flag"]
    return amap


# ---------------------------------------------------------------------------
# Chart builder  (matplotlib → PNG bytes embedded in slide)
# ---------------------------------------------------------------------------

def make_chart(df, col_ios, col_and, ylabel,
               y_min=None, y_max=None, fmt=None) -> bytes:
    fig, ax = plt.subplots(figsize=(7.8, 4.5), dpi=120)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    # FIX 5: x-axis labels in "DD Mon" format
    weeks = [fmt_date_short(w) for w in df["week_start"]]
    ios_v = df[col_ios].tolist()
    and_v = df[col_and].tolist()

    ax.plot(weeks, ios_v, color=M_NAVY,   lw=2.5, marker="o", ms=7, zorder=3)
    ax.plot(weeks, and_v, color=M_ORANGE, lw=2.5, marker="o", ms=7, zorder=3)

    lbl = fmt if fmt else fmt_k
    for w, v in zip(weeks, ios_v):
        ax.annotate(lbl(v), (w, v), xytext=(0, 9),
                    textcoords="offset points", ha="center",
                    fontsize=8, color=M_NAVY, fontweight="bold")
    for w, v in zip(weeks, and_v):
        ax.annotate(lbl(v), (w, v), xytext=(0, -14),
                    textcoords="offset points", ha="center",
                    fontsize=8, color=M_ORANGE, fontweight="bold")

    # FIX 6: Y axis range
    if y_min is not None and y_max is not None:
        ax.set_ylim(y_min, y_max)
    elif y_min is not None:
        ax.set_ylim(bottom=y_min)

    ax.set_ylabel(ylabel, fontsize=9, color=M_DGRAY)
    ax.tick_params(labelsize=8, colors=M_DGRAY)
    ax.grid(axis="y", color=M_GRID, lw=0.7, zorder=0)
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color(M_GRID)
    ax.spines["bottom"].set_color(M_GRID)
    # FIX 5: horizontal tick labels, no rotation
    plt.xticks(rotation=0, ha="center", fontsize=8)

    handles = [
        Line2D([0], [0], color=M_NAVY,   lw=2.5, marker="o", ms=7, label="● iOS"),
        Line2D([0], [0], color=M_ORANGE, lw=2.5, marker="o", ms=7, label="● Android"),
    ]
    ax.legend(handles=handles, loc="upper center", bbox_to_anchor=(0.5, 1.12),
              ncol=2, frameon=False, fontsize=10)

    plt.tight_layout(pad=1.2)
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", facecolor="white", dpi=120)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# Slide 1 — Cover
# ---------------------------------------------------------------------------

def slide_cover(prs: Presentation, week: str):
    slide = new_slide(prs)
    set_slide_bg(slide, LGRAY)

    sp = slide.shapes.add_shape(9, Inches(8.0), Inches(-1.2), Inches(7.0), Inches(10.0))
    sp.fill.solid()
    sp.fill.fore_color.rgb = TEAL
    sp.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(5.5), Inches(0.4))
    run = txb.text_frame.paragraphs[0].add_run()
    run.text = "SunExpress Airlines"
    run.font.size = Pt(12); run.font.color.rgb = NAVY; run.font.name = "Calibri"

    txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(7.2), Inches(1.1))
    run2 = txb2.text_frame.paragraphs[0].add_run()
    run2.text = "SunExpress"
    run2.font.size = Pt(48); run2.font.bold = True
    run2.font.color.rgb = NAVY; run2.font.name = "Calibri"

    txb3 = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(7.2), Inches(0.9))
    run3 = txb3.text_frame.paragraphs[0].add_run()
    run3.text = "Performance Overview"
    run3.font.size = Pt(36); run3.font.color.rgb = NAVY; run3.font.name = "Calibri"

    txb4 = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(7.2), Inches(0.6))
    run4 = txb4.text_frame.paragraphs[0].add_run()
    run4.text = "12 Jan – 08 Mar 2026"
    run4.font.size = Pt(18); run4.font.color.rgb = ORANGE; run4.font.name = "Calibri"

    txb5 = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(2.0), Inches(0.4))
    run5 = txb5.text_frame.paragraphs[0].add_run()
    run5.text = "loodos"
    run5.font.size = Pt(12); run5.font.bold = True
    run5.font.color.rgb = TEAL; run5.font.name = "Calibri"

    return slide


# ---------------------------------------------------------------------------
# Slides 2 & 3 — Platform Overview
# ---------------------------------------------------------------------------

def slide_overview(prs: Presentation, df, platform: str, anomaly_map: dict):
    slide = new_slide(prs)
    set_slide_bg(slide, LGRAY)

    plat_label = "iOS" if platform == "ios" else "Android"
    cols = (
        ["ios_rating", "ios_crash_free_user", "ios_downloads",
         "ios_active_users", "ios_uninstalls"]
        if platform == "ios" else
        ["android_rating", "android_crash_free_user", "android_downloads",
         "android_active_users", "android_uninstalls"]
    )

    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.75))
    tf  = txb.text_frame; tf.word_wrap = False
    p   = tf.paragraphs[0]
    r1  = p.add_run()
    r1.text = plat_label; r1.font.bold = True; r1.font.color.rgb = ORANGE
    r1.font.size = Pt(36); r1.font.name = "Calibri"
    r2  = p.add_run()
    r2.text = " Performance Overview"; r2.font.bold = True; r2.font.color.rgb = NAVY
    r2.font.size = Pt(36); r2.font.name = "Calibri"

    # FIX 3: column names match reference PDF
    headers    = ["Date", "App Rating", "Crash Free User",
                  "Download (Weekly)", "Active User (Weekly)", "Uninstall Count (Weekly)"]
    col_widths = [Inches(1.9), Inches(1.6), Inches(1.9),
                  Inches(2.0), Inches(2.3), Inches(3.0)]  # total ≈ 12.7"
    nrows = len(df) + 1

    tbl = slide.shapes.add_table(
        nrows, 6, Inches(0.3), Inches(0.95), Inches(12.7), Inches(3.5)
    ).table
    style_header_row(tbl, headers, col_widths)

    for ri, (_, row) in enumerate(df.iterrows()):
        dr = ri + 1
        tbl.rows[dr].height = Inches(0.35)   # FIX 4: compact rows
        # FIX 1: date range format
        style_data_cell(tbl.cell(dr, 0),
                        fmt_date_range(row["week_start"], row["week_end"]))
        for ci, col in enumerate(cols):
            flag = anomaly_map.get((row["week_start"], col))
            style_data_cell(tbl.cell(dr, ci + 1), fmt_metric(row[col], col), flag=flag)

    add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Slides 4–8 — Metric detail  (chart + comparison table)
# ---------------------------------------------------------------------------

def slide_metric(prs: Presentation, df, title: str,
                 col_ios: str, col_and: str, ylabel: str,
                 anomaly_map: dict,
                 y_min=None, y_max=None, fmt=None):
    slide = new_slide(prs)
    set_slide_bg(slide, LGRAY)
    add_title(slide, [(title, True, NAVY, 36)])

    png = make_chart(df, col_ios, col_and, ylabel, y_min, y_max, fmt)
    slide.shapes.add_picture(io.BytesIO(png),
                              Inches(0.3), Inches(0.9), Inches(7.8), Inches(5.8))

    headers    = ["Date", "iOS", "Android"]
    col_widths = [Inches(1.6), Inches(1.5), Inches(1.5)]  # total = 4.6"
    nrows = len(df) + 1

    tbl = slide.shapes.add_table(
        nrows, 3, Inches(8.4), Inches(0.9), Inches(4.6), Inches(3.5)
    ).table
    style_header_row(tbl, headers, col_widths)

    for ri, (_, row) in enumerate(df.iterrows()):
        dr = ri + 1
        tbl.rows[dr].height = Inches(0.35)   # FIX 4: compact rows
        # FIX 1: short date for narrow column
        style_data_cell(tbl.cell(dr, 0), fmt_date_short(row["week_start"]))
        flag_i = anomaly_map.get((row["week_start"], col_ios))
        flag_a = anomaly_map.get((row["week_start"], col_and))
        style_data_cell(tbl.cell(dr, 1), fmt_metric(row[col_ios], col_ios), flag=flag_i)
        style_data_cell(tbl.cell(dr, 2), fmt_metric(row[col_and], col_and), flag=flag_a)

    add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Slide 9 — Trend Summary & Insights
# ---------------------------------------------------------------------------

def slide_trends(prs: Presentation, trends_data: dict, insights_data: dict):
    slide = new_slide(prs)
    set_slide_bg(slide, LGRAY)
    add_title(slide, [("Trend Summary & Insights", True, NAVY, 32)])

    METRIC_KEYS = [
        ("rating",          "App Rating"),
        ("crash_free_user", "Crash Free"),
        ("downloads",       "Downloads"),
        ("active_users",    "Active Users"),
        ("uninstalls",      "Uninstalls"),
    ]
    trend_map = trends_data.get("trends", {})

    def icon(d):   return "▲" if d == "increasing" else ("▼" if d == "decreasing" else "→")
    def tcolor(d): return GREEN if d == "increasing" else (RED if d == "decreasing" else SGRAY)

    tbl = slide.shapes.add_table(
        6, 3, Inches(0.3), Inches(0.95), Inches(5.8), Inches(3.5)
    ).table
    tbl.columns[0].width = Inches(2.1)
    tbl.columns[1].width = Inches(1.85)
    tbl.columns[2].width = Inches(1.85)
    style_header_row(tbl, ["Metric", "iOS", "Android"])

    for ri, (key, label) in enumerate(METRIC_KEYS):
        dr = ri + 1
        tbl.rows[dr].height = Inches(0.52)
        ios_t = trend_map.get(f"ios_{key}", {})
        and_t = trend_map.get(f"android_{key}", {})

        cell = tbl.cell(dr, 0)
        cell.text = label
        run  = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(12); run.font.bold = True; run.font.name = "Calibri"
        set_cell_bg(cell, WHITE)

        for t_data, col_idx in [(ios_t, 1), (and_t, 2)]:
            c = tbl.cell(dr, col_idx)
            c.text_frame.clear()
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            direction = t_data.get("direction", "stable")
            pct       = t_data.get("total_change_pct", 0)
            r.text = f"{icon(direction)} {pct:+.1f}%"
            r.font.size = Pt(12); r.font.bold = True
            r.font.color.rgb = tcolor(direction); r.font.name = "Calibri"
            set_cell_bg(c, WHITE)

    txb_h = slide.shapes.add_textbox(Inches(6.4), Inches(0.95), Inches(6.6), Inches(0.45))
    rh = txb_h.text_frame.paragraphs[0].add_run()
    rh.text = "Key Insights"; rh.font.bold = True; rh.font.size = Pt(16)
    rh.font.color.rgb = NAVY; rh.font.name = "Calibri"

    PRIO_COLORS = {"critical": RED, "important": ORANGE, "info": NAVY}
    PRIO_LABELS = {"critical": "CRITICAL", "important": "IMPORTANT", "info": "INFO"}

    y = 1.48
    for ins in insights_data.get("insights", [])[:5]:
        prio = ins.get("priority", "info")
        obs  = ins.get("observation", "")
        act  = ins.get("action", "")
        if len(obs) > 90: obs = obs[:87] + "..."
        if len(act) > 78: act = act[:75] + "..."

        txb = slide.shapes.add_textbox(Inches(6.4), Inches(y), Inches(6.6), Inches(0.88))
        tf  = txb.text_frame; tf.word_wrap = True

        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = f"[{PRIO_LABELS.get(prio, 'INFO')}] "
        r1.font.bold = True; r1.font.size = Pt(10)
        r1.font.color.rgb = PRIO_COLORS.get(prio, NAVY); r1.font.name = "Calibri"
        r2 = p1.add_run()
        r2.text = obs; r2.font.size = Pt(9)
        r2.font.color.rgb = DGRAY; r2.font.name = "Calibri"

        p2 = tf.add_paragraph()
        r3 = p2.add_run()
        r3.text = f"→ {act}"
        r3.font.size = Pt(8); r3.font.italic = True
        r3.font.color.rgb = SGRAY; r3.font.name = "Calibri"

        y += 1.05

    add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="SunExpress PPTX rapor üretici")
    parser.add_argument("--week",   default="2026-03-02", help="Hafta etiketi YYYY-MM-DD")
    parser.add_argument("--output", default=None,         help="Çıktı dosyası adı (opsiyonel)")
    args = parser.parse_args()
    week = args.week

    print(f"\n{'='*60}")
    print(f"  generate-report.py — week: {week}")
    print(f"{'='*60}")

    print("\n[1/3] Veriler yükleniyor...")
    proc = BASE_DIR / "data" / "processed"
    df       = pd.read_csv(proc / f"weekly-metrics-{week}.csv")
    with open(proc / f"insights-{week}.json",      encoding="utf-8") as f: insights  = json.load(f)
    with open(proc / f"anomalies-{week}.json",     encoding="utf-8") as f: anomalies = json.load(f)
    with open(proc / f"trend-summary-{week}.json", encoding="utf-8") as f: trends    = json.load(f)
    amap = build_anomaly_map(anomalies)
    print(f"  ✓ {len(df)} hafta veri,  {len(amap)} anomali flag")

    print("\n[2/3] Slide'lar oluşturuluyor...")
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_cover(prs, week)
    print("  ✓ Slide 1: Kapak")

    slide_overview(prs, df, "ios",     amap)
    print("  ✓ Slide 2: iOS Performance Overview")

    slide_overview(prs, df, "android", amap)
    print("  ✓ Slide 3: Android Performance Overview")

    # FIX 6 & 7: Y-axis ranges + K-format data labels
    slide_metric(prs, df, "App Rating",
                 "ios_rating", "android_rating", "Rating (★)",
                 amap,
                 y_min=0, y_max=4,                        # PDF: 0–4
                 fmt=lambda v: f"{v:.1f}")
    print("  ✓ Slide 4: App Rating")

    slide_metric(prs, df, "Crash Free User",
                 "ios_crash_free_user", "android_crash_free_user", "Crash-Free Rate (%)",
                 amap,
                 y_min=98, y_max=100,                     # PDF: 98–100
                 fmt=lambda v: f"{v:.2f}")
    print("  ✓ Slide 5: Crash Free User")

    slide_metric(prs, df, "Download (Weekly)",
                 "ios_downloads", "android_downloads", "Downloads",
                 amap,
                 y_min=0,                                  # starts at 0
                 fmt=fmt_k)
    print("  ✓ Slide 6: Download (Weekly)")

    slide_metric(prs, df, "Active User (Weekly)",
                 "ios_active_users", "android_active_users", "Active Users",
                 amap,
                 y_min=0,                                  # starts at 0
                 fmt=fmt_k)
    print("  ✓ Slide 7: Active User (Weekly)")

    slide_metric(prs, df, "Uninstall Count (Weekly)",
                 "ios_uninstalls", "android_uninstalls", "Uninstalls",
                 amap,
                 y_min=0,                                  # starts at 0
                 fmt=fmt_k)
    print("  ✓ Slide 8: Uninstall Count (Weekly)")

    slide_trends(prs, trends, insights)
    print("  ✓ Slide 9: Trend Summary & Insights")

    print("\n[3/3] Dosya kaydediliyor...")
    out_dir  = BASE_DIR / "reports" / week
    out_dir.mkdir(parents=True, exist_ok=True)
    fname    = args.output or f"sunexpress-weekly-report-{week}.pptx"
    out_path = out_dir / fname
    prs.save(str(out_path))
    size_kb  = out_path.stat().st_size / 1024
    print(f"  ✓ {out_path.name}  ({size_kb:.1f} KB)")
    print(f"\n  Rapor: {out_path}")
    print(f"{'='*60}\n")
    return str(out_path)


if __name__ == "__main__":
    main()
