"""
render-slides.py  — PPTX slide thumbnail renderer using matplotlib
Generates slide-NN.jpg for each slide in the PPTX.
"""

import sys
import os
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import matplotlib.image as mpimg

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import numpy as np

WEEK = "2026-03-02"
REPORT_DIR = Path(f"reports/{WEEK}")
PPTX_PATH = REPORT_DIR / f"sunexpress-weekly-report-{WEEK}.pptx"

W_IN = 13.333
H_IN = 7.5
DPI = 144  # 1920x1080 roughly

def rgb_to_hex(rgb):
    if rgb is None:
        return "#CCCCCC"
    return f"#{rgb.rgb:06X}"

def shape_text(shape):
    try:
        return shape.text_frame.text
    except Exception:
        return ""

def shape_fill_color(shape):
    try:
        fill = shape.fill
        if fill.type is not None:
            return rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    return None

def render_slide(slide, slide_idx, out_path):
    fig, ax = plt.subplots(figsize=(W_IN, H_IN), dpi=DPI)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_aspect("equal")
    ax.axis("off")

    # Background
    fig.patch.set_facecolor("#F0F0F0")
    ax.set_facecolor("#F0F0F0")

    slide_w = slide.shapes._spTree.getparent().attrib.get(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}cx", None
    )

    # Use EMU coordinates: 1 EMU = 1/914400 inch
    # slide width = 12192000 EMU, height = 6858000 EMU
    EMU_W = 12192000
    EMU_H = 6858000

    rendered_shapes = []
    for shape in slide.shapes:
        try:
            l = shape.left / EMU_W if shape.left is not None else 0
            t = 1 - (shape.top + shape.height) / EMU_H if shape.top is not None else 0
            w = shape.width / EMU_W if shape.width is not None else 0
            h = shape.height / EMU_H if shape.height is not None else 0

            # Draw shape background
            fill_color = None
            try:
                fill = shape.fill
                if hasattr(fill, 'fore_color') and fill.fore_color is not None:
                    fill_color = f"#{fill.fore_color.rgb:06X}"
            except Exception:
                pass

            if fill_color and fill_color not in ("#FFFFFF", "#F0F0F0", "None"):
                rect = FancyBboxPatch((l, t), w, h,
                                      boxstyle="round,pad=0",
                                      facecolor=fill_color,
                                      edgecolor="none",
                                      alpha=0.85,
                                      zorder=2)
                ax.add_patch(rect)
            else:
                rect = FancyBboxPatch((l, t), w, h,
                                      boxstyle="round,pad=0",
                                      facecolor="white",
                                      edgecolor="#CCCCCC",
                                      linewidth=0.3,
                                      alpha=0.5,
                                      zorder=1)
                ax.add_patch(rect)

            # Draw text
            text = shape_text(shape)
            if text and text.strip():
                # Get font size and color from first run
                font_size = 10
                font_color = "#333333"
                try:
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                font_size = run.font.size / 12700  # EMU to pt
                            if run.font.color and run.font.color.type is not None:
                                font_color = f"#{run.font.color.rgb:06X}"
                            break
                        break
                except Exception:
                    pass

                # Scale font size to figure
                display_size = max(5, min(font_size * 0.55, 18))
                # Truncate long text
                display_text = text[:120].replace("\n", " | ")

                ax.text(l + w / 2, t + h / 2,
                        display_text,
                        ha="center", va="center",
                        fontsize=display_size,
                        color=font_color,
                        wrap=True,
                        clip_on=True,
                        zorder=5)

            # Draw embedded images (charts etc.)
            if hasattr(shape, "image"):
                try:
                    img_data = shape.image.blob
                    img_arr = mpimg.imread(io.BytesIO(img_data))
                    ax.imshow(img_arr,
                              extent=[l, l + w, t, t + h],
                              aspect="auto",
                              zorder=3)
                except Exception:
                    # Placeholder for image
                    rect_img = FancyBboxPatch((l, t), w, h,
                                              boxstyle="round,pad=0",
                                              facecolor="#E8E8E8",
                                              edgecolor="#AAAAAA",
                                              linewidth=0.5,
                                              zorder=3)
                    ax.add_patch(rect_img)
                    ax.text(l + w/2, t + h/2, "[Chart Image]",
                            ha="center", va="center", fontsize=8,
                            color="#666666", zorder=6)

            # Draw tables
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = shape.table
                nrows = len(tbl.rows)
                ncols = len(tbl.columns)
                row_h = h / nrows if nrows > 0 else h
                col_w = w / ncols if ncols > 0 else w
                for r in range(nrows):
                    for c in range(ncols):
                        try:
                            cell = tbl.cell(r, c)
                            cx = l + c * col_w
                            cy = t + (nrows - 1 - r) * row_h
                            cell_fill = "#1B3A5C" if r == 0 else "#FFFFFF"
                            try:
                                cf = cell.fill
                                if cf.fore_color is not None:
                                    cell_fill = f"#{cf.fore_color.rgb:06X}"
                            except Exception:
                                pass
                            cell_rect = plt.Rectangle((cx, cy), col_w, row_h,
                                                      facecolor=cell_fill,
                                                      edgecolor="#CCCCCC",
                                                      linewidth=0.3,
                                                      zorder=4)
                            ax.add_patch(cell_rect)
                            cell_text = cell.text_frame.text[:30] if cell.text_frame else ""
                            text_color = "#FFFFFF" if r == 0 else "#333333"
                            if cell_text:
                                ax.text(cx + col_w/2, cy + row_h/2,
                                        cell_text,
                                        ha="center", va="center",
                                        fontsize=5, color=text_color,
                                        clip_on=True, zorder=7)
                        except Exception:
                            pass

        except Exception as e:
            pass

    # Slide number label
    ax.text(0.98, 0.02, f"Slide {slide_idx}",
            ha="right", va="bottom", fontsize=8,
            color="#888888", transform=ax.transAxes, zorder=10)

    plt.tight_layout(pad=0)
    fig.savefig(str(out_path), format="jpeg", dpi=DPI, bbox_inches="tight",
                facecolor="#F0F0F0")
    plt.close(fig)
    print(f"  Saved: {out_path}")


def main():
    if not PPTX_PATH.exists():
        print(f"ERROR: PPTX not found at {PPTX_PATH}")
        sys.exit(1)

    prs = Presentation(str(PPTX_PATH))
    REPORT_DIR.mkdir(parents=True, exist_ok=True)

    print(f"Rendering {len(prs.slides)} slides from {PPTX_PATH}")
    for i, slide in enumerate(prs.slides, start=1):
        out_path = REPORT_DIR / f"slide-{i:02d}.jpg"
        render_slide(slide, i, out_path)

    print(f"\nDone. {len(prs.slides)} thumbnails saved to {REPORT_DIR}/")


if __name__ == "__main__":
    main()
