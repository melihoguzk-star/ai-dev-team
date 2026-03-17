"""
qa-report.py — Visual & Data QA for SunExpress weekly PPTX report
Outputs results to reports/WEEK/qa-visual-report.md
"""

import sys
import json
import csv
import io
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

WEEK = "2026-03-02"
REPORT_DIR = Path(f"reports/{WEEK}")
PPTX_PATH = REPORT_DIR / f"sunexpress-weekly-report-{WEEK}.pptx"
CSV_PATH = Path(f"data/processed/weekly-metrics-{WEEK}.csv")

# Brand colors
NAVY   = "1B3A5C"
ORANGE = "E85D26"
TEAL   = "2ABFBF"
LGRAY  = "F0F0F0"
WHITE  = "FFFFFF"
WARN   = "FFF3CD"
CRIT   = "F8D7DA"


# ── helpers ──────────────────────────────────────────────────────────────────

def rgb_hex(rgb) -> str:
    """Convert RGBColor to hex string — handles python-pptx 3.14 format bug."""
    try:
        # Try the standard way first
        return f"{rgb.rgb:06X}"
    except Exception:
        pass
    try:
        # Fallback: access underlying bytes directly
        v = rgb._RGBColor__int  # private attribute in some versions
        return f"{v:06X}"
    except Exception:
        pass
    try:
        return str(rgb)
    except Exception:
        return ""


def xml_srgb(element, *paths) -> str:
    """Walk XML paths to find a:srgbClr and return its 'val' attribute (uppercase)."""
    # paths: list of qname sequences to try
    for path in paths:
        cur = element
        for tag in path:
            if cur is None:
                break
            cur = cur.find(tag)
        if cur is not None:
            val = cur.get("val")
            if val:
                return val.upper()
    return ""


def shape_fill_hex_xml(shape) -> str:
    """Extract fill color from shape XML (works around python-pptx RGBColor bug)."""
    try:
        sp = shape._element
        # spPr > a:solidFill > a:srgbClr
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            spPr = sp.find(qn("p:grpSpPr"))
        if spPr is not None:
            sf = spPr.find(qn("a:solidFill"))
            if sf is not None:
                srgb = sf.find(qn("a:srgbClr"))
                if srgb is not None:
                    return srgb.get("val", "").upper()
    except Exception:
        pass
    return ""


def run_color_xml(run) -> str:
    """Extract font color from run XML."""
    try:
        r_el = run._r
        rPr = r_el.find(qn("a:rPr"))
        if rPr is not None:
            sf = rPr.find(qn("a:solidFill"))
            if sf is not None:
                srgb = sf.find(qn("a:srgbClr"))
                if srgb is not None:
                    return srgb.get("val", "").upper()
    except Exception:
        pass
    # Fallback to high-level API
    try:
        return rgb_hex(run.font.color.rgb)
    except Exception:
        return ""


def cell_fill_hex_xml(cell) -> str:
    """Extract table cell fill color via XML."""
    try:
        tc = cell._tc
        tcPr = tc.find(qn("a:tcPr"))
        if tcPr is not None:
            sf = tcPr.find(qn("a:solidFill"))
            if sf is not None:
                srgb = sf.find(qn("a:srgbClr"))
                if srgb is not None:
                    return srgb.get("val", "").upper()
    except Exception:
        pass
    return ""


def shape_text(shape) -> str:
    try:
        return shape.text_frame.text
    except Exception:
        return ""

def all_texts(slide) -> list[str]:
    return [shape_text(s) for s in slide.shapes]

def slide_has_text(slide, fragment: str) -> bool:
    return any(fragment.lower() in t.lower() for t in all_texts(slide))

def slide_has_image(slide) -> bool:
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            return True
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    return False

def get_tables(slide):
    tables = []
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                _ = s.table  # verify it's actually a table
                tables.append(s.table)
            except Exception:
                pass
    return tables

def fill_hex(cell):
    """Use XML-based cell fill detection."""
    return cell_fill_hex_xml(cell)

def slide_bg_color(slide):
    """Try to determine slide background fill color hex via XML."""
    try:
        bg = slide.background
        cSld = bg._element
        # Try via fill object first
        fill = bg.fill
        try:
            return rgb_hex(fill.fore_color.rgb)
        except Exception:
            pass
        # Try XML path: p:bg > p:bgPr > a:solidFill > a:srgbClr
        bgPr = cSld.find(qn("p:bgPr"))
        if bgPr is not None:
            sf = bgPr.find(qn("a:solidFill"))
            if sf is not None:
                srgb = sf.find(qn("a:srgbClr"))
                if srgb is not None:
                    return srgb.get("val", "").upper()
    except Exception:
        pass
    return ""

def table_dimensions(table):
    return len(table.rows), len(table.columns)

def cell_text(tbl, row, col) -> str:
    try:
        return tbl.cell(row, col).text_frame.text.strip()
    except Exception:
        return ""


# ── per-slide checks ──────────────────────────────────────────────────────────

def check_slide1(slide):
    results = {}
    texts = all_texts(slide)

    results["has_sunexpress_text"] = any("SunExpress" in t for t in texts)
    results["has_performance_overview"] = any("Performance Overview" in t for t in texts)
    results["has_date_range"] = any("2026" in t for t in texts)
    results["has_loodos_footer"] = any("loodos" in t.lower() for t in texts)

    # Check decorative shape (non-text shape with fill) — use XML detection
    has_shape = False
    for shape in slide.shapes:
        color = shape_fill_hex_xml(shape)
        if color.upper() == TEAL:
            has_shape = True
            break
    results["has_decorative_shape"] = has_shape

    # Background check
    bg = slide_bg_color(slide)
    results["background_is_lgray"] = bg.upper() == LGRAY or bg == ""

    # Title color check (48pt bold navy) — use XML detection
    has_navy_title = False
    for shape in slide.shapes:
        try:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if "SunExpress" in run.text:
                        color = run_color_xml(run)
                        if color.upper() == NAVY:
                            has_navy_title = True
        except Exception:
            pass
    results["title_is_navy"] = has_navy_title

    return results


def check_slide_platform(slide, platform: str):
    results = {}
    tables = get_tables(slide)
    texts = all_texts(slide)

    results["has_platform_name_in_title"] = any(platform in t for t in texts)
    results["has_performance_overview"] = any("Performance Overview" in t for t in texts)
    results["has_loodos_footer"] = any("loodos" in t.lower() for t in texts)
    results["has_sunexpress_airlines_footer"] = any("SunExpress Airlines" in t for t in texts)

    if tables:
        tbl = tables[0]
        nrows, ncols = table_dimensions(tbl)
        results["table_rows"] = nrows
        results["table_cols"] = ncols
        results["table_has_8_rows"] = nrows == 8
        results["table_has_6_cols"] = ncols == 6

        # Check all 7 weeks present in Date column
        date_col_idx = 0
        dates_found = []
        for r in range(1, nrows):
            ct = cell_text(tbl, r, date_col_idx)
            if ct:
                dates_found.append(ct)
        results["weeks_in_table"] = len(dates_found)
        results["all_7_weeks_present"] = len(dates_found) >= 7

        # Header background color — XML-based detection
        try:
            hdr_color = cell_fill_hex_xml(tbl.cell(0, 0))
            results["header_color"] = hdr_color.upper()
            results["header_is_navy"] = hdr_color.upper() == NAVY
        except Exception:
            results["header_is_navy"] = False
    else:
        results["table_found"] = False

    return results


def check_slide_metric(slide, metric_name: str):
    results = {}
    texts = all_texts(slide)

    results["has_metric_title"] = any(metric_name in t for t in texts)
    results["has_image"] = slide_has_image(slide)
    results["has_loodos_footer"] = any("loodos" in t.lower() for t in texts)
    results["has_sunexpress_airlines_footer"] = any("SunExpress Airlines" in t for t in texts)

    tables = get_tables(slide)
    if tables:
        tbl = tables[0]
        nrows, ncols = table_dimensions(tbl)
        results["table_rows"] = nrows
        results["table_cols"] = ncols
        results["table_has_8_rows"] = nrows == 8
        results["table_has_3_cols"] = ncols == 3
    else:
        results["table_found"] = False

    return results


def check_slide9(slide):
    results = {}
    texts = all_texts(slide)
    tables = get_tables(slide)

    results["has_loodos_footer"] = any("loodos" in t.lower() for t in texts)
    results["has_sunexpress_airlines_footer"] = any("SunExpress Airlines" in t for t in texts)

    if tables:
        tbl = tables[0]
        nrows, ncols = table_dimensions(tbl)
        results["trend_table_rows"] = nrows
        results["trend_table_cols"] = ncols
        results["trend_table_ok"] = nrows >= 6 and ncols >= 3
    else:
        results["trend_table_found"] = False

    # Count insight text boxes
    insight_count = sum(1 for t in texts if len(t.strip()) > 30)
    results["insight_textboxes_count"] = insight_count
    results["has_3_or_more_insights"] = insight_count >= 3

    return results


# ── data cross-check ─────────────────────────────────────────────────────────

def parse_number(s: str) -> float | None:
    """Parse formatted number string to float."""
    try:
        s = s.strip().replace(",", "").replace("%", "").replace(" ", "")
        if s == "" or s == "-":
            return None
        return float(s)
    except Exception:
        return None

def cross_check_data(prs):
    results = {"ios": [], "android": [], "matched": 0, "total": 0, "score_pct": 0}

    # Read CSV
    csv_data = []
    with open(CSV_PATH, newline="") as f:
        reader = csv.DictReader(f)
        for row in reader:
            csv_data.append(row)

    ios_metrics = ["ios_rating", "ios_crash_free_user", "ios_downloads", "ios_active_users", "ios_uninstalls"]
    android_metrics = ["android_rating", "android_crash_free_user", "android_downloads", "android_active_users", "android_uninstalls"]

    def fmt(val_str):
        try:
            v = float(val_str)
            if v >= 1000:
                return f"{int(v):,}"
            elif v > 10:
                return f"{v:.2f}"
            else:
                return f"{v:.2f}"
        except Exception:
            return str(val_str)

    # Slide 2 — iOS
    slide2 = prs.slides[1]
    tables2 = get_tables(slide2)
    if tables2:
        tbl = tables2[0]
        matched = 0
        total = 0
        for r_idx, row in enumerate(csv_data):
            tbl_row = r_idx + 1
            if tbl_row >= len(tbl.rows):
                break
            for c_idx, metric in enumerate(ios_metrics):
                col = c_idx + 1
                if col >= len(tbl.columns):
                    continue
                csv_val = parse_number(row[metric])
                tbl_val = parse_number(cell_text(tbl, tbl_row, col))
                total += 1
                if csv_val is not None and tbl_val is not None:
                    tol = max(abs(csv_val) * 0.01, 0.05)
                    if abs(csv_val - tbl_val) <= tol:
                        matched += 1
                        results["ios"].append({
                            "week": row["week_start"], "metric": metric,
                            "csv": csv_val, "pptx": tbl_val, "match": True
                        })
                    else:
                        results["ios"].append({
                            "week": row["week_start"], "metric": metric,
                            "csv": csv_val, "pptx": tbl_val, "match": False
                        })
        results["ios_matched"] = matched
        results["ios_total"] = total

    # Slide 3 — Android
    slide3 = prs.slides[2]
    tables3 = get_tables(slide3)
    if tables3:
        tbl = tables3[0]
        matched = 0
        total = 0
        for r_idx, row in enumerate(csv_data):
            tbl_row = r_idx + 1
            if tbl_row >= len(tbl.rows):
                break
            for c_idx, metric in enumerate(android_metrics):
                col = c_idx + 1
                if col >= len(tbl.columns):
                    continue
                csv_val = parse_number(row[metric])
                tbl_val = parse_number(cell_text(tbl, tbl_row, col))
                total += 1
                if csv_val is not None and tbl_val is not None:
                    tol = max(abs(csv_val) * 0.01, 0.05)
                    if abs(csv_val - tbl_val) <= tol:
                        matched += 1
                        results["android"].append({
                            "week": row["week_start"], "metric": metric,
                            "csv": csv_val, "pptx": tbl_val, "match": True
                        })
                    else:
                        results["android"].append({
                            "week": row["week_start"], "metric": metric,
                            "csv": csv_val, "pptx": tbl_val, "match": False
                        })
        results["android_matched"] = matched
        results["android_total"] = total

    total_matched = results.get("ios_matched", 0) + results.get("android_matched", 0)
    total_all = results.get("ios_total", 0) + results.get("android_total", 0)
    results["matched"] = total_matched
    results["total"] = total_all
    results["score_pct"] = round(total_matched / total_all * 100, 1) if total_all > 0 else 0

    return results


# ── branding compliance ───────────────────────────────────────────────────────

def check_branding(prs):
    score = 0
    max_score = 0
    details = []

    def check(name, passed, weight=10):
        nonlocal score, max_score
        max_score += weight
        if passed:
            score += weight
        details.append({"check": name, "passed": passed, "weight": weight})

    slide1 = prs.slides[0]
    bg = slide_bg_color(slide1)
    check("Kapak arka planı LGRAY (#F0F0F0)", bg.upper() == LGRAY or bg == "", weight=10)

    # NAVY header check across slides — use XML-based cell fill detection
    navy_header_ok = True
    for s_idx in [1, 2]:  # slides 2,3
        sl = prs.slides[s_idx]
        for tbl in get_tables(sl):
            try:
                hdr = cell_fill_hex_xml(tbl.cell(0, 0))
                if hdr.upper() != NAVY:
                    navy_header_ok = False
            except Exception:
                pass
    check("Tablo header rengi NAVY (#1B3A5C)", navy_header_ok, weight=10)

    # Orange for platform names — use XML run color detection
    orange_found = False
    for s_idx in [1, 2]:
        sl = prs.slides[s_idx]
        for shape in sl.shapes:
            try:
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.text in ("iOS", "Android"):
                            c = run_color_xml(run)
                            if c.upper() == ORANGE:
                                orange_found = True
            except Exception:
                pass
    check("Platform adı turuncu (#E85D26)", orange_found, weight=10)

    # Teal loodos footer — use XML run color detection
    teal_found = False
    for sl in prs.slides:
        for shape in sl.shapes:
            try:
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        if "loodos" in run.text.lower():
                            c = run_color_xml(run)
                            if c.upper() == TEAL:
                                teal_found = True
            except Exception:
                pass
    check("Loodos footer teal (#2ABFBF)", teal_found, weight=10)

    # Critical cells (#F8D7DA) — use XML cell fill detection
    crit_found = False
    for sl in prs.slides:
        for tbl in get_tables(sl):
            for r in range(len(tbl.rows)):
                for c in range(len(tbl.columns)):
                    ch = cell_fill_hex_xml(tbl.cell(r, c))
                    if ch.upper() == CRIT:
                        crit_found = True
    check("Kritik hücre rengi (#F8D7DA)", crit_found, weight=10)

    # Warning cells (#FFF3CD) — use XML cell fill detection
    warn_found = False
    for sl in prs.slides:
        for tbl in get_tables(sl):
            for r in range(len(tbl.rows)):
                for c in range(len(tbl.columns)):
                    ch = cell_fill_hex_xml(tbl.cell(r, c))
                    if ch.upper() == WARN:
                        warn_found = True
    check("Uyarı hücre rengi (#FFF3CD)", warn_found, weight=10)

    # All slides have loodos footer
    all_have_loodos = all(
        any("loodos" in t.lower() for t in all_texts(sl))
        for sl in prs.slides
    )
    check("Tüm slide'larda loodos footer var", all_have_loodos, weight=10)

    # All slides have SunExpress Airlines footer
    all_have_se = all(
        any("SunExpress Airlines" in t for t in all_texts(sl))
        for sl in prs.slides
    )
    check("Tüm slide'larda SunExpress Airlines footer var", all_have_se, weight=10)

    # Slide dimensions 16:9
    prs_w = prs.slide_width
    prs_h = prs.slide_height
    ratio = prs_w / prs_h if prs_h else 0
    check("Slide boyutu 16:9 (13.33\" x 7.5\")", abs(ratio - 16/9) < 0.01, weight=10)

    # 9 slides total
    check("Toplam 9 slide", len(prs.slides) == 9, weight=10)

    pct = round(score / max_score * 100, 1) if max_score > 0 else 0
    return {"score": score, "max_score": max_score, "pct": pct, "details": details}


# ── main QA runner ────────────────────────────────────────────────────────────

def run_qa():
    report_lines = []
    append = report_lines.append

    prs = Presentation(str(PPTX_PATH))

    append("# SunExpress Haftalık Rapor — Görsel QA Sonucu")
    append("")
    append(f"**Rapor Haftası:** {WEEK}")
    append(f"**QA Tarihi:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    append(f"**PPTX Dosyası:** `{PPTX_PATH}`")
    append("")

    # Generation status
    pptx_exists = PPTX_PATH.exists()
    pptx_size = PPTX_PATH.stat().st_size / 1024 if pptx_exists else 0
    append("## 1. Üretim Durumu")
    append("")
    append(f"| Kontrol | Sonuç |")
    append(f"|---------|-------|")
    append(f"| PPTX oluşturuldu | {'✅ PASS' if pptx_exists else '❌ FAIL'} |")
    append(f"| PPTX boyutu | {pptx_size:.1f} KB |")
    append(f"| Slide sayısı | {len(prs.slides)} |")
    append("")

    # PDF status
    pdf_path = REPORT_DIR / f"sunexpress-weekly-report-{WEEK}.pdf"
    pdf_exists = pdf_path.exists()
    append("## 2. PDF Dönüşüm Durumu")
    append("")
    append(f"| Kontrol | Sonuç |")
    append(f"|---------|-------|")
    append(f"| PDF oluşturuldu | {'✅ PASS' if pdf_exists else '⚠️ HAYIR'} |")
    if not pdf_exists:
        append(f"| Sebep | LibreOffice yok; pptx2pdf ve aspose-slides Python 3.14 ile uyumsuz |")
        append(f"| Etki | PDF çıktısı mevcut değil — PPTX manuel dönüşüm gerektirebilir |")
    append("")

    # JPG status
    jpg_files = sorted(REPORT_DIR.glob("slide-*.jpg"))
    append("## 3. JPG Thumbnail Durumu")
    append("")
    append(f"| Kontrol | Sonuç |")
    append(f"|---------|-------|")
    append(f"| Thumbnail üretildi | {'✅ PASS' if len(jpg_files) == 9 else '❌ FAIL'} |")
    append(f"| Dosya sayısı | {len(jpg_files)}/9 |")
    for jpg in jpg_files:
        size_kb = jpg.stat().st_size / 1024
        append(f"| {jpg.name} | {size_kb:.1f} KB |")
    append("")

    # Per-slide QA
    append("## 4. Slide Bazlı QA Sonuçları")
    append("")

    all_slide_results = {}

    # Slide 1
    append("### Slide 1 — Kapak")
    s1 = check_slide1(prs.slides[0])
    all_slide_results[1] = s1
    append("")
    append("| Kontrol | Sonuç |")
    append("|---------|-------|")
    for k, v in s1.items():
        label_map = {
            "has_sunexpress_text": "'SunExpress' metni var",
            "has_performance_overview": "'Performance Overview' metni var",
            "has_date_range": "Tarih aralığı metni var",
            "has_loodos_footer": "'loodos' footer var",
            "has_decorative_shape": "Dekoratif şekil var (teal)",
            "background_is_lgray": "Arka plan LGRAY (#F0F0F0)",
            "title_is_navy": "Başlık rengi NAVY (#1B3A5C)",
        }
        label = label_map.get(k, k)
        status = "✅ PASS" if v else "❌ FAIL"
        append(f"| {label} | {status} |")
    append("")

    # Slides 2-3
    for s_idx, platform in [(2, "iOS"), (3, "Android")]:
        append(f"### Slide {s_idx} — {platform} Performance Overview")
        sr = check_slide_platform(prs.slides[s_idx - 1], platform)
        all_slide_results[s_idx] = sr
        append("")
        append("| Kontrol | Sonuç |")
        append("|---------|-------|")
        label_map2 = {
            "has_platform_name_in_title": f"'{platform}' başlıkta var",
            "has_performance_overview": "'Performance Overview' var",
            "has_loodos_footer": "'loodos' footer var",
            "has_sunexpress_airlines_footer": "'SunExpress Airlines' footer var",
            "table_has_8_rows": "Tablo 8 satır (1 header + 7 veri)",
            "table_has_6_cols": "Tablo 6 sütun",
            "all_7_weeks_present": "Tabloda 7 hafta var",
            "header_is_navy": "Header rengi NAVY",
        }
        for k, v in sr.items():
            if k in label_map2:
                label = label_map2[k]
                if isinstance(v, bool):
                    status = "✅ PASS" if v else "❌ FAIL"
                else:
                    status = str(v)
                append(f"| {label} | {status} |")
        append("")

    # Slides 4-8
    metric_slides = [
        (4, "App Rating"),
        (5, "Crash Free User"),
        (6, "Download"),
        (7, "Active User"),
        (8, "Uninstall"),
    ]
    for s_num, metric in metric_slides:
        append(f"### Slide {s_num} — {metric}")
        sr = check_slide_metric(prs.slides[s_num - 1], metric)
        all_slide_results[s_num] = sr
        append("")
        append("| Kontrol | Sonuç |")
        append("|---------|-------|")
        label_map3 = {
            "has_metric_title": f"'{metric}' başlıkta var",
            "has_image": "Grafik PNG embed var",
            "has_loodos_footer": "'loodos' footer var",
            "has_sunexpress_airlines_footer": "'SunExpress Airlines' footer var",
            "table_has_8_rows": "Tablo 8 satır",
            "table_has_3_cols": "Tablo 3 sütun (Date | iOS | Android)",
        }
        for k, v in sr.items():
            if k in label_map3:
                label = label_map3[k]
                if isinstance(v, bool):
                    status = "✅ PASS" if v else "❌ FAIL"
                else:
                    status = str(v)
                append(f"| {label} | {status} |")
        append("")

    # Slide 9
    append("### Slide 9 — Trend Summary & Key Insights")
    s9 = check_slide9(prs.slides[8])
    all_slide_results[9] = s9
    append("")
    append("| Kontrol | Sonuç |")
    append("|---------|-------|")
    label_map9 = {
        "has_loodos_footer": "'loodos' footer var",
        "has_sunexpress_airlines_footer": "'SunExpress Airlines' footer var",
        "trend_table_ok": "Trend tablosu (6+ satır × 3+ sütun)",
        "has_3_or_more_insights": "En az 3 insight metin kutusu",
    }
    for k, v in s9.items():
        if k in label_map9:
            label = label_map9[k]
            if isinstance(v, bool):
                status = "✅ PASS" if v else "❌ FAIL"
            else:
                status = str(v)
            append(f"| {label} | {status} |")
    append(f"| Insight metin sayısı | {s9.get('insight_textboxes_count', 0)} |")
    append(f"| Trend tablo satır × sütun | {s9.get('trend_table_rows', '?')} × {s9.get('trend_table_cols', '?')} |")
    append("")

    # Data cross-check
    append("## 5. Veri Doğruluğu (CSV ↔ PPTX)")
    append("")
    data_results = cross_check_data(prs)
    ios_matched = data_results.get("ios_matched", 0)
    ios_total = data_results.get("ios_total", 0)
    and_matched = data_results.get("android_matched", 0)
    and_total = data_results.get("android_total", 0)
    total_matched = data_results["matched"]
    total_all = data_results["total"]
    score_pct = data_results["score_pct"]

    append(f"| Platform | Eşleşen / Toplam | Oran |")
    append(f"|----------|-----------------|------|")
    ios_pct = round(ios_matched/ios_total*100, 1) if ios_total else 0
    and_pct = round(and_matched/and_total*100, 1) if and_total else 0
    append(f"| iOS (Slide 2) | {ios_matched}/{ios_total} | {ios_pct}% |")
    append(f"| Android (Slide 3) | {and_matched}/{and_total} | {and_pct}% |")
    append(f"| **Toplam** | **{total_matched}/{total_all}** | **{score_pct}%** |")
    append("")

    # Show mismatches
    mismatches = [x for x in data_results.get("ios", []) + data_results.get("android", []) if not x["match"]]
    if mismatches:
        append("### Uyumsuz Değerler")
        append("")
        append("| Hafta | Metrik | CSV Değeri | PPTX Değeri |")
        append("|-------|--------|-----------|------------|")
        for mm in mismatches:
            append(f"| {mm['week']} | {mm['metric']} | {mm['csv']} | {mm['pptx']} |")
        append("")
    else:
        append("Tüm değerler eşleşiyor. ✅")
        append("")

    # Branding compliance
    append("## 6. Branding Uyum Kontrolü")
    append("")
    brand = check_branding(prs)
    append(f"**Branding Skoru: {brand['score']}/{brand['max_score']} ({brand['pct']}%)**")
    append("")
    append("| Kontrol | Ağırlık | Sonuç |")
    append("|---------|---------|-------|")
    for d in brand["details"]:
        status = "✅ PASS" if d["passed"] else "❌ FAIL"
        append(f"| {d['check']} | {d['weight']} | {status} |")
    append("")

    # Issues found
    append("## 7. Tespit Edilen Sorunlar")
    append("")
    issues = []

    if not pdf_exists:
        issues.append("PDF dönüşümü yapılamadı — LibreOffice, pptx2pdf ve aspose-slides Python 3.14 ile uyumsuz.")

    for k, v in s1.items():
        if isinstance(v, bool) and not v:
            issues.append(f"Slide 1: '{k}' kontrolü başarısız.")

    for s_idx in [2, 3]:
        sr = all_slide_results[s_idx]
        for k, v in sr.items():
            if isinstance(v, bool) and not v:
                issues.append(f"Slide {s_idx}: '{k}' kontrolü başarısız.")

    for s_num in range(4, 9):
        sr = all_slide_results[s_num]
        for k, v in sr.items():
            if isinstance(v, bool) and not v:
                issues.append(f"Slide {s_num}: '{k}' kontrolü başarısız.")

    for k, v in s9.items():
        if isinstance(v, bool) and not v:
            issues.append(f"Slide 9: '{k}' kontrolü başarısız.")

    for d in brand["details"]:
        if not d["passed"]:
            issues.append(f"Branding: '{d['check']}' başarısız.")

    if score_pct < 90:
        issues.append(f"Veri doğruluğu %{score_pct} — bazı değerler CSV ile uyuşmuyor.")

    if issues:
        for i, issue in enumerate(issues, 1):
            append(f"{i}. {issue}")
    else:
        append("Hiçbir sorun tespit edilmedi. ✅")
    append("")

    # Overall score
    append("## 8. Genel Skor")
    append("")

    # Count pass/fail
    total_checks = 0
    passed_checks = 0

    for sr in all_slide_results.values():
        for k, v in sr.items():
            if isinstance(v, bool):
                total_checks += 1
                if v:
                    passed_checks += 1

    # Add branding checks
    for d in brand["details"]:
        total_checks += 1
        if d["passed"]:
            passed_checks += 1

    # Add data accuracy
    total_checks += 1
    if score_pct >= 90:
        passed_checks += 1

    # Add generation checks
    total_checks += 2
    if pptx_exists: passed_checks += 1
    if len(jpg_files) == 9: passed_checks += 1

    overall_pct = round(passed_checks / total_checks * 100, 1) if total_checks else 0

    grade = "MÜKEMMEL" if overall_pct >= 95 else \
            "BAŞARILI" if overall_pct >= 80 else \
            "KABUL EDİLEBİLİR" if overall_pct >= 70 else "BAŞARISIZ"

    append(f"| Bileşen | Skor |")
    append(f"|---------|------|")
    append(f"| Slide QA (pass/fail) | {passed_checks - (1 if score_pct >= 90 else 0) - (1 if pptx_exists else 0) - (1 if len(jpg_files)==9 else 0)}/{total_checks - 3} |")
    append(f"| Veri Doğruluğu | {score_pct}% |")
    append(f"| Branding Uyumu | {brand['pct']}% |")
    append(f"| PPTX Üretimi | {'✅' if pptx_exists else '❌'} |")
    append(f"| JPG Thumbnail | {'✅ 9/9' if len(jpg_files)==9 else '❌'} |")
    append(f"| PDF Dönüşümü | {'✅' if pdf_exists else '⚠️ Yapılamadı'} |")
    append(f"| **Genel Kontrol Oranı** | **{passed_checks}/{total_checks} ({overall_pct}%)** |")
    append(f"| **Sonuç** | **{grade}** |")
    append("")
    append("---")
    append(f"*Bu rapor `qa-report.py` tarafından otomatik oluşturulmuştur — {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*")
    append(f"*Loodos Reporting Pipeline v1.0*")

    # Write to file
    out_path = REPORT_DIR / "qa-visual-report.md"
    out_path.write_text("\n".join(report_lines), encoding="utf-8")
    print(f"\nQA raporu: {out_path}")

    # Print summary to stdout
    print(f"\n{'='*60}")
    print(f"  QA ÖZET — Hafta {WEEK}")
    print(f"{'='*60}")
    print(f"  PPTX üretildi:     {'✓' if pptx_exists else '✗'}")
    print(f"  PDF dönüşüm:       {'✓' if pdf_exists else '✗ (kütüphane yok)'}")
    print(f"  JPG thumbnail:     {len(jpg_files)}/9")
    print(f"  Veri doğruluğu:    {total_matched}/{total_all} ({score_pct}%)")
    print(f"  Branding skoru:    {brand['score']}/{brand['max_score']} ({brand['pct']}%)")
    print(f"  Genel skor:        {passed_checks}/{total_checks} ({overall_pct}%)")
    print(f"  Sonuç:             {grade}")
    print(f"{'='*60}")

    return issues, overall_pct


if __name__ == "__main__":
    issues, pct = run_qa()
    sys.exit(0 if pct >= 70 else 1)
